from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Crear presentación
prs = Presentation()

# --- SLIDE 1: Portada ---
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Automatización de Procesamiento y Gestión de Datos - LogiGol App"
subtitle.text = "Proyecto desarrollado por los emprendedores de la futura empresa LogiGol Systems\n© 2025 LogiGol Sky"

# --- SLIDE 2: Objetivo del Proyecto ---
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
slide_2.shapes.title.text = "Objetivo del Sistema"
content = slide_2.placeholders[1]
content.text = (
    "Desarrollar una aplicación en Python con PyQt5 que permita automatizar:\n"
    "• La lectura, concatenación y validación de archivos CSV.\n"
    "• La creación y actualización de tablas en una base de datos SQLite.\n"
    "• La exportación de datos a CSV y Excel.\n"
    "• La gestión de información deportiva de forma eficiente y visual."
)

# --- SLIDE 3: Arquitectura General ---
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
slide_3.shapes.title.text = "Arquitectura del Sistema"
content = slide_3.placeholders[1]
content.text = (
    "La solución está compuesta por los siguientes módulos:\n\n"
    "1. Interfaz gráfica construida con PyQt5.\n"
    "2. Motor de procesamiento de datos con pandas y numpy.\n"
    "3. Base de datos SQLite para almacenamiento persistente.\n"
    "4. Módulo de exportación e importación de CSV/Excel.\n"
    "5. Funciones de validación y clasificación de datos deportivos."
)

# --- SLIDE 4: Funcionalidades Implementadas ---
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
slide_4.shapes.title.text = "Funcionalidades Implementadas"
content = slide_4.placeholders[1]
content.text = (
    "• Selección de carpeta con archivos CSV.\n"
    "• Concatenación automática de archivos con columna 'Archivo'.\n"
    "• Creación de tablas CL_O, clasificacion_local y clasificacion_visitante.\n"
    "• Copia y resguardo de base de datos original.\n"
    "• Exportación a CSV desde SQLite.\n"
    "• Guardado de datos procesados en la base de datos copiada."
)

# --- SLIDE 5: Interfaz Gráfica PyQt5 ---
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
slide_5.shapes.title.text = "Interfaz Gráfica con PyQt5"
content = slide_5.placeholders[1]
content.text = (
    "Ventana principal con los siguientes botones:\n\n"
    "• Seleccionar Carpeta\n"
    "• Crear Tablas en BD\n"
    "• Exportar BD a CSV\n"
    "• Concatenar CSV\n"
    "• Copiar Base de Datos\n"
    "• Guardar en BD Copiada\n\n"
    "Cada botón ejecuta una función específica del flujo de procesamiento."
)

# --- SLIDE 6: Solución a Errores Encontrados ---
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
slide_6.shapes.title.text = "Corrección de Errores y Mejoras"
content = slide_6.placeholders[1]
content.text = (
    "• Error `sqlite3.OperationalError: no such table: CL_O` solucionado creando las tablas antes de insertar.\n"
    "• Advertencias de PyQt5 (`Unused argument 'parent'`) ignoradas sin impacto funcional.\n"
    "• Se agregó manejo de excepciones y mensajes visuales con QMessageBox.\n"
    "• Se incluyó exportación automática a CSV y Excel tras concatenar."
)

# --- SLIDE 7: Flujo de Procesamiento ---
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
slide_7.shapes.title.text = "Flujo de Procesamiento de Datos"
content = slide_7.placeholders[1]
content.text = (
    "1. Selección de carpeta con archivos CSV.\n"
    "2. Concatenación y validación de datos.\n"
    "3. Creación de tablas en base de datos copiada.\n"
    "4. Inserción y complementación de datos.\n"
    "5. Exportación final a CSV.\n"
    "6. Visualización en interfaz PyQt5."
)

# --- SLIDE 8: Próximos Pasos ---
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
slide_8.shapes.title.text = "Próximos Pasos y Mejoras Futuras"
content = slide_8.placeholders[1]
content.text = (
    "• Integración de análisis estadístico de rendimiento deportivo.\n"
    "• Implementación de visualizaciones interactivas (matplotlib/plotly).\n"
    "• Conexión a servicios web o API deportivas en tiempo real.\n"
    "• Optimización del almacenamiento y consultas SQL.\n"
    "• Implementación de un sistema multiusuario con autenticación."
)

# --- SLIDE 9: Créditos y Autores ---
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
slide_9.shapes.title.text = "Créditos y Autores"
content = slide_9.placeholders[1]
content.text = (
    "Desarrollado por:\n"
    "👨‍💻 Equipo de Emprendedores de LogiGol Systems\n\n"
    "Colaboradores Técnicos:\n"
    "• Ingenieros en Datos y Software\n"
    "• Diseñadores de Interfaz UX/UI\n\n"
    "Empresa: LogiGol Sky - Innovación en Gestión de Datos Deportivos\n"
    "Año: 2025"
)

# Guardar presentación
output_path = "/mnt/data/Presentacion_LogiGol_Automatizacion.pptx"
prs.save(output_path)
output_path
